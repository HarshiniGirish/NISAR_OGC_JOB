from __future__ import annotations

from pathlib import Path
import textwrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "User_Code_to_App_Package_Demo_Presentation.pptx"
ASSET_DIR = ROOT / "presentation_assets"
BRANCH = "cursor/ai-ogcification-system-ad8b"
GITHUB = f"https://github.com/HarshiniGirish/NISAR_OGC_JOB/blob/{BRANCH}"

COLORS = {
    "navy": RGBColor(20, 44, 80),
    "blue": RGBColor(38, 102, 169),
    "light_blue": RGBColor(226, 240, 255),
    "green": RGBColor(48, 134, 83),
    "orange": RGBColor(214, 117, 35),
    "gray": RGBColor(90, 96, 105),
    "light_gray": RGBColor(245, 247, 250),
    "white": RGBColor(255, 255, 255),
}


def font(size: int = 20, bold: bool = False):
    try:
        return ImageFont.truetype(
            "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf" if not bold else "/usr/share/fonts/truetype/dejavu/DejaVuSansMono-Bold.ttf",
            size,
        )
    except OSError:
        return ImageFont.load_default()


def read_lines(path: str, start: int, end: int) -> str:
    lines = (ROOT / path).read_text(encoding="utf-8").splitlines()
    selected = lines[start - 1 : end]
    return "\n".join(f"{idx + start:>4}  {line}" for idx, line in enumerate(selected))


def make_code_image(name: str, title: str, code: str, width: int = 1500) -> Path:
    ASSET_DIR.mkdir(exist_ok=True)
    wrapped = []
    for line in code.splitlines():
        if len(line) <= 120:
            wrapped.append(line)
        else:
            wrapped.extend(textwrap.wrap(line, width=120, subsequent_indent="      "))
    line_height = 28
    height = 90 + max(1, len(wrapped)) * line_height + 35
    img = Image.new("RGB", (width, height), (32, 36, 44))
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, width, 58), fill=(22, 82, 140))
    draw.text((26, 16), title, fill=(255, 255, 255), font=font(23, bold=True))
    y = 80
    for line in wrapped:
        draw.text((26, y), line, fill=(232, 238, 245), font=font(20))
        y += line_height
    path = ASSET_DIR / name
    img.save(path)
    return path


def make_architecture_image() -> Path:
    ASSET_DIR.mkdir(exist_ok=True)
    img = Image.new("RGB", (1700, 950), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    boxes = [
        ("Notebook / script", 90, 60),
        ("Parser + static analysis", 90, 170),
        ("DPS-readiness scan", 90, 280),
        ("MCP defaults + dataset facts", 90, 390),
        ("LLM recommendations / V2 synthesis", 90, 500),
        ("Dependency graph + access planner", 90, 610),
        ("Package generator", 90, 720),
        ("CWL + OGC AP validator", 90, 830),
    ]
    right_boxes = [
        ("Generated artifacts", 940, 225),
        ("run.sh / build.sh / env.yml", 940, 335),
        ("application.cwl / workflow.cwl", 940, 445),
        ("application-package.cwl", 940, 555),
        ("readiness + final reports", 940, 665),
    ]
    for text, x, y in boxes:
        draw.rounded_rectangle((x, y, x + 650, y + 70), radius=18, fill=(226, 240, 255), outline=(38, 102, 169), width=3)
        draw.text((x + 22, y + 20), text, fill=(20, 44, 80), font=font(27, bold=True))
    for (_, x, y), (_, _x2, y2) in zip(boxes, boxes[1:]):
        draw.line((415, y + 70, 415, y2), fill=(38, 102, 169), width=5)
        draw.polygon([(405, y2 - 8), (425, y2 - 8), (415, y2 + 10)], fill=(38, 102, 169))
    for text, x, y in right_boxes:
        draw.rounded_rectangle((x, y, x + 650, y + 70), radius=18, fill=(240, 247, 241), outline=(48, 134, 83), width=3)
        draw.text((x + 22, y + 20), text, fill=(25, 98, 62), font=font(25, bold=True))
    draw.line((740, 755, 940, 590), fill=(214, 117, 35), width=6)
    draw.polygon([(930, 585), (952, 582), (941, 604)], fill=(214, 117, 35))
    draw.text((860, 150), "Reviewable outputs", fill=(214, 117, 35), font=font(32, bold=True))
    path = ASSET_DIR / "architecture.png"
    img.save(path)
    return path


def make_dependency_image() -> Path:
    ASSET_DIR.mkdir(exist_ok=True)
    img = Image.new("RGB", (1650, 850), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    imports = ["earthaccess", "h5py", "s3fs", "xarray", "maap", "rasterio"]
    packages = ["earthaccess", "h5py + h5netcdf", "s3fs + fsspec + boto3", "xarray + netcdf4 + scipy", "maap-py", "rasterio"]
    draw.text((60, 40), "Dependency graph explains provenance", fill=(20, 44, 80), font=font(34, bold=True))
    for i, imp in enumerate(imports):
        y = 130 + i * 105
        draw.rounded_rectangle((70, y, 430, y + 65), radius=14, fill=(226, 240, 255), outline=(38, 102, 169), width=2)
        draw.text((92, y + 18), f"import {imp}", fill=(20, 44, 80), font=font(22, bold=True))
        draw.rounded_rectangle((900, y, 1500, y + 65), radius=14, fill=(240, 247, 241), outline=(48, 134, 83), width=2)
        draw.text((922, y + 18), packages[i], fill=(25, 98, 62), font=font(22, bold=True))
        draw.line((430, y + 32, 900, y + 32), fill=(90, 96, 105), width=4)
        draw.polygon([(885, y + 22), (905, y + 32), (885, y + 42)], fill=(90, 96, 105))
    draw.text((500, 760), "Manual map + implicit rules + access-plan deps; LLM suggestions must validate first", fill=(90, 96, 105), font=font(22))
    path = ASSET_DIR / "dependency_graph.png"
    img.save(path)
    return path


def add_title(slide, title: str, subtitle: str = ""):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["white"]
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = COLORS["navy"]
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12), Inches(0.38))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = COLORS["gray"]


def add_bullets(slide, bullets, x=0.75, y=1.45, w=6.1, h=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["navy"]
        p.level = 0


def add_image(slide, image_path: Path, x, y, w=None, h=None):
    kwargs = {}
    if w:
        kwargs["width"] = Inches(w)
    if h:
        kwargs["height"] = Inches(h)
    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), **kwargs)


def add_link_box(slide, label: str, url: str, y: float):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(y), Inches(12), Inches(0.45))
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["light_blue"]
    box.line.color.rgb = COLORS["blue"]
    tx = box.text_frame.paragraphs[0]
    tx.text = f"{label}: {url}"
    tx.font.size = Pt(12)
    tx.font.color.rgb = COLORS["blue"]


def add_footer(slide, n: int):
    box = slide.shapes.add_textbox(Inches(11.8), Inches(7.05), Inches(1.1), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = str(n)
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["gray"]
    p.alignment = PP_ALIGN.RIGHT


def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide_no = 1

    arch = make_architecture_image()
    dep = make_dependency_image()
    scanner_img = make_code_image(
        "dps_scanner_code.png",
        "generator/dps_readiness_scan.py",
        read_lines("generator/dps_readiness_scan.py", 12, 35),
    )
    resolver_img = make_code_image(
        "mcp_resolver_code.png",
        "mcp_server/tools/default_resolver.py",
        read_lines("mcp_server/tools/default_resolver.py", 20, 55),
    )
    llm_img = make_code_image(
        "llm_v2_code.png",
        "generator/suggested_notebook_v2.py",
        read_lines("generator/suggested_notebook_v2.py", 30, 78),
    )
    validator_img = make_code_image(
        "validator_code.png",
        "generator/ogc_validator.py",
        read_lines("generator/ogc_validator.py", 64, 100),
    )
    command_img = make_code_image(
        "demo_commands.png",
        "Live demo commands",
        textwrap.dedent(
            """
            python3 generator/generate_package.py \\
              --input notebooks/gedi_calval_maap_stack_ogcification_demo.ipynb \\
              --scan-dps-readiness --use-mcp-defaults --llm-recommendations \\
              --build-dependency-graph --emit-suggested-notebook-v2 \\
              --validate-ogc --output-dir generated_gedi_calval_demo

            cd generated_gedi_calval_demo
            cwltool --validate application.cwl
            cwltool --validate workflow.cwl
            ap-validator --format json --detail all application-package.cwl
            python3 gedi_calval_maap_stack_ogcification_demo.py
            """
        ).strip(),
    )

    # 1
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["navy"]
    title = slide.shapes.add_textbox(Inches(0.75), Inches(1.1), Inches(12), Inches(1.5))
    p = title.text_frame.paragraphs[0]
    p.text = "AI-assisted User Code → MAAP DPS / OGC Application Package"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS["white"]
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.75), Inches(11.7), Inches(1))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "NISAR_OGC_JOB demonstration: notebooks, MCP metadata, LLM review, dependency graphing, V2 notebook synthesis, and OGC validation"
    sp.font.size = Pt(20)
    sp.font.color.rgb = COLORS["white"]
    add_footer(slide, slide_no); slide_no += 1

    # 2
    slide = prs.slides.add_slide(blank); add_title(slide, "Problem: notebooks are not automatically deployable")
    add_bullets(slide, [
        "Science workflows often start as exploratory notebooks with hidden cell state.",
        "MAAP DPS jobs must be non-interactive, parameterized, reproducible, and output to output/.",
        "OGC Application Packages require CWL, containers, metadata, inputs/outputs, and validation.",
        "Manual conversion is repetitive and support-intensive: wrappers, envs, dependencies, credentials, access strategy, validation.",
    ], w=11.8, size=22)
    add_footer(slide, slide_no); slide_no += 1

    # 3
    slide = prs.slides.add_slide(blank); add_title(slide, "Implemented end-to-end architecture")
    add_image(slide, arch, 0.45, 1.15, w=12.35)
    add_footer(slide, slide_no); slide_no += 1

    # 4
    slide = prs.slides.add_slide(blank); add_title(slide, "Repository map: what was implemented")
    add_bullets(slide, [
        "generator/: parser, scanner, dependency graph, access planner, validator, reports, notebook V2.",
        "mcp_server/tools/: CMR, asset inspection, access options, subset-cost, default resolver.",
        "templates/: Dockerfile, run.sh, env.yml, algorithm.yml, application.cwl, workflow.cwl.",
        "notebooks/: MAAP-docs-style sample and GEDI Cal/Val MAAP STAC demo notebook.",
        "tests/: 35 unit/integration tests for generator, MCP tools, access planning, and OGCification pipeline.",
    ], w=11.8, size=20)
    add_link_box(slide, "Repository", "https://github.com/HarshiniGirish/NISAR_OGC_JOB", 6.65)
    add_footer(slide, slide_no); slide_no += 1

    # 5
    slide = prs.slides.add_slide(blank); add_title(slide, "Demo inputs: MAAP docs + MAAP STAC dataset")
    add_bullets(slide, [
        "NISAR Access tutorial: direct S3 credentials + Earthaccess HTTPS fallback.",
        "GEDI Cal/Val demo: uses GEDI_CalVal_Lidar_Data from the MAAP STAC stack.",
        "The demo notebook intentionally contains notebook-only visualization and DPS-ready metadata processing sections.",
        "The pipeline shows original notebook → readiness scan → suggested V2 → generated package → validation → run outputs.",
    ], w=11.8, size=20)
    add_link_box(slide, "MAAP NISAR docs", "https://docs.maap-project.org/en/develop/science/NISAR/NISAR_access.html", 5.8)
    add_link_box(slide, "GEDI Cal/Val STAC", "https://stac-browser.maap-project.org/collections/GEDI_CalVal_Lidar_Data", 6.35)
    add_footer(slide, slide_no); slide_no += 1

    # 6
    slide = prs.slides.add_slide(blank); add_title(slide, "DPS-readiness scanner")
    add_bullets(slide, [
        "Classifies notebook cells/functions/blocks as DPS-ready, candidate-after-refactor, notebook-only, or blocking.",
        "Detects imports, data access, outputs, hardcoded values, plotting, interactive code, local paths, and dependencies.",
        "Produces dps_readiness_report.json and dps_readiness_report.md.",
    ], x=0.65, y=1.3, w=4.4, size=18)
    add_image(slide, scanner_img, 5.0, 1.2, w=7.8)
    add_link_box(slide, "Code", f"{GITHUB}/generator/dps_readiness_scan.py", 6.8)
    add_footer(slide, slide_no); slide_no += 1

    # 7
    slide = prs.slides.add_slide(blank); add_title(slide, "Example NISAR readiness results")
    add_bullets(slide, [
        "DPS-ready cells: S3 credentials, Earthaccess discovery, S3 open path, HTTPS open path.",
        "Notebook-only cells: markdown, display, matplotlib/folium visualization, inspection cells.",
        "Info warnings: implicit notebook state; these are not failures but highlight ordering/refactor risks.",
        "Outcome: V2 notebook makes runtime values explicit and packages the DPS-relevant path.",
    ], w=11.8, size=21)
    add_footer(slide, slide_no); slide_no += 1

    # 8
    slide = prs.slides.add_slide(blank); add_title(slide, "MCP metadata/default resolver")
    add_bullets(slide, [
        "Creates provenance-tagged default suggestions for collection_id, short_name, asset_href, bbox, variables, access_mode, output_directory.",
        "Does not hardcode MAAP docs values as real defaults.",
        "Sources: user manifest, source evidence, MAAP/STAC/CMR facts, safe fallback logic.",
    ], x=0.65, y=1.3, w=4.35, size=18)
    add_image(slide, resolver_img, 5.0, 1.2, w=7.8)
    add_link_box(slide, "Code", f"{GITHUB}/mcp_server/tools/default_resolver.py", 6.8)
    add_footer(slide, slide_no); slide_no += 1

    # 9
    slide = prs.slides.add_slide(blank); add_title(slide, "Dependency graph: explain why every dependency exists")
    add_image(slide, dep, 0.55, 1.15, w=12.25)
    add_link_box(slide, "Code", f"{GITHUB}/generator/dependency_graph.py", 6.9)
    add_footer(slide, slide_no); slide_no += 1

    # 10
    slide = prs.slides.add_slide(blank); add_title(slide, "Access strategy planner")
    add_bullets(slide, [
        "Rule-based planner always works; OpenAI planner is optional.",
        "Allowed strategies include direct_s3_xarray, direct_s3_h5py, rasterio_windowed_read, zarr_open_zarr, harmony_subset, cmr_search_then_s3, stac_raster_api, HTTPS fallback.",
        "AI cannot invent unsupported strategies; every plan is validated.",
        "Writes access_plan.json and access_runtime.py with implementation hints.",
    ], w=11.9, size=20)
    add_link_box(slide, "Code", f"{GITHUB}/generator/access_planner.py", 6.65)
    add_footer(slide, slide_no); slide_no += 1

    # 11
    slide = prs.slides.add_slide(blank); add_title(slide, "LLM guardrails + suggested notebook V2")
    add_bullets(slide, [
        "LLM is optional: core generation remains deterministic.",
        "OpenAI can synthesize a runtime-safe V2 notebook from hidden-state/readiness analysis.",
        "Returned notebook JSON is schema-checked and code cells must compile.",
        "Original notebook is never overwritten; fallback V2 is emitted if LLM fails.",
    ], x=0.65, y=1.3, w=4.35, size=18)
    add_image(slide, llm_img, 5.0, 1.2, w=7.8)
    add_link_box(slide, "Code", f"{GITHUB}/generator/suggested_notebook_v2.py", 6.8)
    add_footer(slide, slide_no); slide_no += 1

    # 12
    slide = prs.slides.add_slide(blank); add_title(slide, "Generated OGC / MAAP package")
    add_bullets(slide, [
        "MAAP DPS artifacts: algorithm.yml, algorithm_config.yaml, run.sh, build.sh, env.yml, Dockerfile, register_dps.py.",
        "OGC artifacts: application.cwl, workflow.cwl, packed application-package.cwl, STAC input/output manifests, publish_ogc.py.",
        "Analysis artifacts: analysis.json, dps_readiness_report.md, dependency_graph.json, dataset_facts.json, access_plan.json, final_readiness_report.json.",
    ], w=11.9, size=20)
    add_footer(slide, slide_no); slide_no += 1

    # 13
    slide = prs.slides.add_slide(blank); add_title(slide, "Validation: cwltool + ap-validator")
    add_bullets(slide, [
        "cwltool validates application.cwl and workflow.cwl syntax.",
        "cwltool --pack creates application-package.cwl.",
        "ap-validator checks OGC EO Application Package best-practice compliance.",
        "ogc_validation_report.md reports OGC ready, MAAP DPS ready, CWL valid, OGC AP valid, warnings, and next actions.",
    ], x=0.65, y=1.3, w=4.35, size=18)
    add_image(slide, validator_img, 5.0, 1.2, w=7.8)
    add_link_box(slide, "Code", f"{GITHUB}/generator/ogc_validator.py", 6.8)
    add_footer(slide, slide_no); slide_no += 1

    # 14
    slide = prs.slides.add_slide(blank); add_title(slide, "GEDI MAAP STAC demo: complete path")
    add_bullets(slide, [
        "Dataset: GEDI_CalVal_Lidar_Data from MAAP STAC.",
        "Notebook recommends an algorithm from a small catalog: metadata inventory, asset manifest, full LAS metrics, rasterization.",
        "DPS-ready path summarizes STAC item metadata and writes output/gedi_calval_summary.json + output/asset_manifest.csv.",
        "V2 notebook can be packaged and run as its own generated OGC/MAAP package.",
    ], w=11.8, size=20)
    add_link_box(slide, "Notebook", f"{GITHUB}/notebooks/gedi_calval_maap_stack_ogcification_demo.ipynb", 6.65)
    add_footer(slide, slide_no); slide_no += 1

    # 15
    slide = prs.slides.add_slide(blank); add_title(slide, "Live demonstration commands")
    add_image(slide, command_img, 0.8, 1.25, w=11.85)
    add_footer(slide, slide_no); slide_no += 1

    # 16
    slide = prs.slides.add_slide(blank); add_title(slide, "What to show live")
    add_bullets(slide, [
        "Open original notebook: show exploratory cells + visualization.",
        "Run generator with readiness/MCP/dependency/LLM/V2/validation flags.",
        "Open dps_readiness_report.md: show DPS-ready vs notebook-only cells.",
        "Open suggested_notebook_v2.ipynb: show parameter cell, preserved helpers, output-writing path.",
        "Validate application-package.cwl with ap-validator: valid: true.",
        "Run generated script and show output files.",
    ], w=11.8, size=21)
    add_footer(slide, slide_no); slide_no += 1

    # 17
    slide = prs.slides.add_slide(blank); add_title(slide, "Important nuance: package validation vs data runtime")
    add_bullets(slide, [
        "OGC-ready means package structure and CWL/Application Package metadata validate.",
        "Runtime may still depend on live services: Earthdata, MAAP UAT, ASF S3 credentials, STAC/CMR endpoints.",
        "Example: MAAP UAT 503 is an access-service issue, not an OGC package failure.",
        "Reports separate structural readiness from runtime access risks.",
    ], w=11.8, size=22)
    add_footer(slide, slide_no); slide_no += 1

    # 18
    slide = prs.slides.add_slide(blank); add_title(slide, "Current limitations and next improvements")
    add_bullets(slide, [
        "Automatic S3 → HTTPS fallback orchestration should be strengthened.",
        "More dataset-specific output templates can improve scientific runtime correctness.",
        "Live MCP metadata integrations can improve defaults and subset-cost estimates.",
        "MAAP DPS sandbox submission/polling should be added for end-to-end cloud job proof.",
        "CI should run ap-validator and representative notebook generation tests.",
    ], w=11.8, size=21)
    add_footer(slide, slide_no); slide_no += 1

    # 19
    slide = prs.slides.add_slide(blank); add_title(slide, "Takeaway")
    add_bullets(slide, [
        "The repository now demonstrates an AI-assisted OGCification workflow.",
        "It explains what is DPS-ready, what stays notebook-only, which dependencies are required, and why.",
        "It generates MAAP DPS and OGC package artifacts and validates the packed OGC Application Package.",
        "It preserves the original notebook and emits a suggested V2 notebook for review or packaging.",
        "The system turns exploratory notebook work into a documented, reviewable, reproducible package candidate.",
    ], w=11.8, size=22)
    add_footer(slide, slide_no); slide_no += 1

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    create_deck()
